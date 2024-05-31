import tkinter as tk
from tkinter import messagebox
import os
import subprocess
from threading import Thread
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image, ImageTk
import fitz  # PyMuPDF

# Configurar o caminho para o executável do Tesseract OCR
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

current_index = 0
processed_files = []

def mostrar_preview(root, arquivo):
    if not os.path.isfile(arquivo):
        return

    preview_window = tk.Toplevel(root)
    preview_window.title("Pré-visualização")
    preview_window.geometry("400x400")
    center_window(preview_window, 400, 400)

    try:
        if arquivo.lower().endswith('.pdf'):
            pdf = fitz.open(arquivo)
            page = pdf.load_page(0)
            pix = page.get_pixmap()
            imagem = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            imagem.thumbnail((400, 400))
            img = ImageTk.PhotoImage(imagem)
            panel = tk.Label(preview_window, image=img)
            panel.image = img
            panel.pack()
        elif arquivo.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
            imagem = Image.open(arquivo)
            imagem.thumbnail((400, 400))
            img = ImageTk.PhotoImage(imagem)
            panel = tk.Label(preview_window, image=img)
            panel.image = img
            panel.pack()
        elif arquivo.lower().endswith('.docx'):
            doc = Document(arquivo)
            texto = "\n".join([para.text for para in doc.paragraphs])
            lbl = tk.Label(preview_window, text=texto, wraplength=400)
            lbl.pack()
        else:
            with open(arquivo, 'r', encoding='iso-8859-1') as f:
                texto = f.read()
            lbl = tk.Label(preview_window, text=texto, wraplength=400)
            lbl.pack()
    except Exception as e:
        print(f"Erro ao mostrar pré-visualização: {e}")

def center_window(window, width, height):
    screen_width = window.winfo_screenwidth()
    screen_height = window.winfo_screenheight()
    position_top = int(screen_height / 2 - height / 2)
    position_right = int(screen_width / 2 - width / 2)
    window.geometry(f'{width}x{height}+{position_right}+{position_top}')

def processar_em_lote(arquivos, output_dir, marker='##40', overlap_interval=2, incluir_imagens=False):
    def processar_arquivo(arquivo):
        texto_extraido = extrair_texto(arquivo)
        if texto_extraido is None:
            print(f"Erro ao extrair texto para o arquivo: {arquivo}")
            return
        nome_arquivo_base = os.path.splitext(os.path.basename(arquivo))[0]
        extensao_arquivo = os.path.splitext(arquivo)[1][1:]  # Remove o ponto inicial da extensão
        subdir = os.path.join(output_dir, extensao_arquivo)
        if not os.path.exists(subdir):
            os.makedirs(subdir)
        novo_caminho = os.path.join(subdir, f"{nome_arquivo_base}.txt")
        texto_formatado = format_text_with_markers(texto_extraido, marker=marker, overlap_interval=overlap_interval)
        with open(novo_caminho, 'w', encoding='utf-8') as novo_arquivo:
            novo_arquivo.write(texto_formatado)
        print(f"Arquivo formatado: {novo_caminho}")

    threads = []
    for idx, arquivo in enumerate(arquivos, start=1):
        if incluir_imagens or not arquivo.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
            print(f"Processando arquivo {idx}/{len(arquivos)}: {arquivo}")
            thread = Thread(target=processar_arquivo, args=(arquivo,))
            threads.append(thread)
            thread.start()

    for thread in threads:
        thread.join()

    print(f"Operação concluída. Os arquivos formatados foram salvos em: {output_dir}")
    abrir_explorador_textos_formatados(output_dir)

def abrir_configuracao_formatacao_lote(callback_formatar, callback_extrair, incluir_imagens):
    global config_window
    config_window = tk.Toplevel()
    config_window.title("Configuração de Formatação")
    center_window(config_window, 500, 250)

    tk.Label(config_window, text="Intervalo de sobreposição:").pack(pady=5)
    overlap_entry = tk.Entry(config_window)
    overlap_entry.insert(0, "2")
    overlap_entry.pack(pady=5)

    tk.Label(config_window, text="Marcador:").pack(pady=5)
    marker_entry = tk.Entry(config_window)
    marker_entry.insert(0, "##40")
    marker_entry.pack(pady=5)

    if incluir_imagens is not None:
        frame_checkbox = tk.Frame(config_window)
        frame_checkbox.pack(fill=tk.X, pady=5)

        checkbox_imagens = tk.Checkbutton(frame_checkbox, text="Formatar Também os Arquivos de Imagem", variable=incluir_imagens)
        checkbox_imagens.pack(pady=5)
        checkbox_imagens.pack(expand=True)

    frame_buttons = tk.Frame(config_window)
    frame_buttons.pack(pady=20)

    def confirmar_configuracao():
        overlap = int(overlap_entry.get())
        marker = marker_entry.get()
        config_window.destroy()
        callback_formatar(overlap, marker)

    btn_confirmar = tk.Button(frame_buttons, text="Confirmar", command=confirmar_configuracao)
    btn_confirmar.pack(side=tk.LEFT, padx=10)

    btn_extrair = tk.Button(frame_buttons, text="Extrair Texto Sem Formatar", command=lambda: [config_window.destroy(), callback_extrair()])
    btn_extrair.pack(side=tk.LEFT, padx=10)

from preview import mostrar_preview  # Certifique-se de importar a função corretamente

def abrir_configuracao_formatacao_individual(callback_formatar, callback_extrair, arquivos):
    global config_window, current_index, processed_files

    if not isinstance(arquivos, list) or not arquivos:
        messagebox.showwarning("Aviso", "Nenhum arquivo disponível para processar.")
        return

    arquivo_nome_completo = arquivos[current_index]
    nome_arquivo = os.path.basename(arquivo_nome_completo)  # Extrair apenas o nome do arquivo com extensão

    config_window = tk.Toplevel()
    config_window.title("Configuração de Formatação")
    center_window(config_window, 500, 250)

    frame_top = tk.Frame(config_window)
    frame_top.pack(fill=tk.X, pady=5)

    label_nome_arquivo = tk.Label(frame_top, text=f"Arquivo: {nome_arquivo}")
    label_nome_arquivo.pack(side=tk.LEFT, padx=5)

    tk.Label(config_window, text="Intervalo de sobreposição:").pack(pady=5)
    overlap_entry = tk.Entry(config_window)
    overlap_entry.insert(0, "2")
    overlap_entry.pack(pady=5)

    tk.Label(config_window, text="Marcador:").pack(pady=5)
    marker_entry = tk.Entry(config_window)
    marker_entry.insert(0, "##40")
    marker_entry.pack(pady=5)

    frame_buttons = tk.Frame(config_window)
    frame_buttons.pack(pady=20)

    def confirmar_configuracao():
        global current_index
        overlap = int(overlap_entry.get())
        marker = marker_entry.get()
        config_window.destroy()
        callback_formatar(overlap, marker)

    btn_confirmar = tk.Button(frame_buttons, text="Confirmar", command=confirmar_configuracao)
    btn_confirmar.pack(side=tk.LEFT, padx=10)

    btn_extrair = tk.Button(frame_buttons, text="Extrair Texto Sem Formatar", command=lambda: [config_window.destroy(), callback_extrair(arquivo_nome_completo)])
    btn_extrair.pack(side=tk.LEFT, padx=10)

    def navegar(delta):
        global current_index
        new_index = current_index + delta
        if 0 <= new_index < len(arquivos) and arquivos[new_index] not in processed_files:
            current_index = new_index
            config_window.destroy()
            abrir_configuracao_formatacao_individual(callback_formatar, callback_extrair, arquivos)
        else:
            messagebox.showwarning("Aviso", "Não é possível acessar este arquivo (lembre-se de selecionar mais de 1 arquivo).")

    btn_voltar = tk.Button(frame_buttons, text="⬅️", command=lambda: navegar(-1))
    btn_voltar.pack(side=tk.LEFT, padx=5)

    btn_avancar = tk.Button(frame_buttons, text="➡️", command=lambda: navegar(1))
    btn_avancar.pack(side=tk.LEFT, padx=5)

    btn_preview = tk.Button(frame_buttons, text="👁", command=lambda: mostrar_preview(config_window, arquivo_nome_completo))
    btn_preview.pack(side=tk.LEFT, padx=5)

def mostrar_mensagem_finalizacao(root):
    def on_continuar():
        messagebox.showinfo("Info", "Voltando para a tela inicial.")
        root.deiconify()
        msg_box.destroy()

    def on_fechar():
        root.quit()

    msg_box = tk.Toplevel(root)
    msg_box.title("Processamento Concluído")
    center_window(msg_box, 360, 140)

    lbl_mensagem = tk.Label(msg_box, text="Processamento concluído. O que você gostaria de fazer?")
    lbl_mensagem.pack(pady=10)

    btn_continuar = tk.Button(msg_box, text="Processar/Extrair mais arquivos", command=on_continuar)
    btn_continuar.pack(side=tk.LEFT, padx=20, pady=20)

    btn_fechar = tk.Button(msg_box, text="Fechar Programa", command=on_fechar)
    btn_fechar.pack(side=tk.RIGHT, padx=20, pady=20)

def processar_arquivo(arquivo, output_dir, marker, overlap_interval):
    texto_extraido = extrair_texto(arquivo)
    if texto_extraido is None:
        print(f"Erro ao extrair texto para o arquivo: {arquivo}")
        return
    nome_arquivo_base = os.path.splitext(os.path.basename(arquivo))[0]
    extensao_arquivo = os.path.splitext(arquivo)[1][1:]  # Remove o ponto inicial da extensão
    subdir = os.path.join(output_dir, extensao_arquivo)
    if not os.path.exists(subdir):
        os.makedirs(subdir)
    novo_caminho = os.path.join(subdir, f"{nome_arquivo_base}.txt")
    texto_formatado = format_text_with_markers(texto_extraido, marker=marker, overlap_interval=overlap_interval)
    with open(novo_caminho, 'w', encoding='utf-8') as novo_arquivo:
        novo_arquivo.write(texto_formatado)
    print(f"Arquivo formatado: {novo_caminho}")

def extrair_textos_sem_formatar(arquivos, output_dir):
    for arquivo in arquivos:
        texto_extraido = extrair_texto(arquivo)
        if texto_extraido is not None:
            nome_arquivo_base = os.path.splitext(os.path.basename(arquivo))[0]
            extensao_arquivo = os.path.splitext(arquivo)[1][1:]  # Remove o ponto inicial da extensão
            subdir = os.path.join(output_dir, extensao_arquivo)
            if not os.path.exists(subdir):
                os.makedirs(subdir)
            novo_caminho = os.path.join(subdir, f"{nome_arquivo_base}.txt")
            with open(novo_caminho, 'w', encoding='utf-8') as novo_arquivo:
                novo_arquivo.write(texto_extraido)
            print(f"Texto extraído salvo em: {novo_caminho}")
    abrir_explorador_textos_formatados(output_dir)

def extrair_texto(arquivo):
    if arquivo.lower().endswith('.pdf'):
        return extrair_texto_pdf(arquivo)
    elif arquivo.lower().endswith('.docx'):
        return extrair_texto_docx(arquivo)
    elif arquivo.lower().endswith('.pptx'):
        return extrair_texto_pptx(arquivo)
    elif arquivo.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
        return extrair_texto_imagem(arquivo)
    elif arquivo.lower().endswith('.txt'):
        with open(arquivo, 'r', encoding='iso-8859-1') as txt_file:
            return txt_file.read()
    else:
        print(f"Formato de arquivo não suportado para: {arquivo}")
        return None

def extrair_texto_pdf_escaneado(caminho_pdf):
    documento = fitz.open(caminho_pdf)
    texto_completo = ''
    for pagina in documento:
        pix = pagina.get_pixmap()
        imagem = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        texto = pytesseract.image_to_string(imagem)
        texto_completo += texto + '\n'
    return texto_completo

def extrair_texto_pdf(caminho_pdf):
    try:
        texto = ''
        with open(caminho_pdf, 'rb') as file:
            pdf_reader = PdfReader(file)
            if not any(page.extract_text() for page in pdf_reader.pages):
                texto = extrair_texto_pdf_escaneado(caminho_pdf)
            else:
                texto = extract_text_from_searchable_pdf(caminho_pdf)
        return texto
    except Exception as e:
        print(f"Erro ao extrair texto do PDF: {e}")
        return None

def extract_text_from_searchable_pdf(caminho_pdf):
    with open(caminho_pdf, 'rb') as file:
        pdf = PdfReader(file)
        texto = ''
        for page in pdf.pages:
            texto += page.extract_text() or ""
        return texto

def extrair_texto_docx(caminho_arquivo):
    doc = Document(caminho_arquivo)
    texto = ""
    for para in doc.paragraphs:
        texto += para.text + "\n"
    return texto

def extrair_texto_pptx(caminho_arquivo):
    prs = Presentation(caminho_arquivo)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texto += shape.text + "\n"
    return texto

def extrair_texto_imagem(caminho_arquivo):
    try:
        texto = pytesseract.image_to_string(Image.open(caminho_arquivo))
        return texto
    except Exception as e:
        print(f"Erro ao extrair texto da imagem: {e}")
        return None

def format_text_with_markers(text, marker='##40', overlap_interval=2):
    lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
    formatted_lines = []
    for i in range(0, len(lines), overlap_interval):
        block_lines = lines[i:i + overlap_interval]
        if block_lines:
            formatted_lines.append(marker)
            formatted_lines.extend(block_lines)
            formatted_lines.append('')
    formatted_text = '\n'.join(formatted_lines)
    return formatted_text

def abrir_explorador_textos_formatados(diretorio):
    subprocess.Popen(f'explorer "{os.path.abspath(diretorio)}"')

# Exemplo de como chamar a função no código maior
def iniciar_processo_de_formatacao(arquivos):
    global current_index
    current_index = 0  # Inicializa o índice antes de iniciar o processo
    abrir_configuracao_formatacao_individual(callback_formatar, callback_extrair, arquivos)

# Exemplo de uso
if __name__ == "__main__":
    root = tk.Tk()
    root.withdraw()  # Esconde a janela principal

    arquivos = ["C:\\caminho\\para\\arquivo1.txt", "C:\\caminho\\para\\arquivo2.txt", "C:\\caminho\\para\\arquivo3.txt"]  # Caminho completo para cada arquivo
    iniciar_processo_de_formatacao(arquivos)

    root.mainloop()